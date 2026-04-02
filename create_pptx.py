#!/usr/bin/env python3
"""
Generate a PowerPoint deck for the App Gateway WAF -> Azure Firewall -> Web App demo.
Covers two use cases:
  1. X-Forwarded-For preservation through App Gateway → Firewall → PE chain
  2. Geofencing WAF + VPN: compound Block rule so VPN bypasses geo-block but NOT managed rules
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────
AZURE_BLUE   = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE    = RGBColor(0x00, 0x2B, 0x5C)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
GREEN        = RGBColor(0x10, 0x7C, 0x10)
RED          = RGBColor(0xD1, 0x34, 0x38)
ORANGE       = RGBColor(0xFF, 0x8C, 0x00)
BLACK        = RGBColor(0x00, 0x00, 0x00)
GRAY_TEXT    = RGBColor(0x60, 0x60, 0x60)
PURPLE       = RGBColor(0x6B, 0x3F, 0xA0)

# No deployment-specific values — all references are generic
TOTAL_SLIDES = 17

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide_content(slide, items, left=1.0, top=2.2, width=11.0, font_size=18, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(10)
        p.level = 0
    return txBox

def add_shape_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Segoe UI"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZURE_BLUE
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_slide_title(slide, title, subtitle=None):
    add_text_box(slide, 0.8, 0.4, 11.5, 0.8, title, font_size=32, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text_box(slide, 0.8, 1.2, 11.5, 0.5, subtitle, font_size=16, color=GRAY_TEXT)

def add_slide_number(slide, num, total=TOTAL_SLIDES):
    add_text_box(slide, 12.0, 7.0, 1.2, 0.4, f"{num}/{total}", font_size=10, color=GRAY_TEXT,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text_box(slide, 1.5, 1.5, 10.0, 1.5,
             "Azure Application Gateway (WAF)\n→ Azure Firewall → Web App",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.2, 10.0, 0.8,
             "X-Forwarded-For Preservation + Geo-Filtering with VPN Security",
             font_size=22, color=RGBColor(0x80, 0xBF, 0xEE), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.5, 10.0, 0.5,
             "Architecture Deep Dive & Live Demo with Proof",
             font_size=18, color=RGBColor(0xA0, 0xA0, 0xA0), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.2, 10.0, 0.5,
             "April 2026",
             font_size=14, color=RGBColor(0x80, 0x80, 0x80), alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Problem Statement", "Why do we need this architecture?")

items = [
    "🔒  Web App must NOT be accessible directly from the Internet",
    "🛡️  All traffic must be inspected by Azure Firewall (L3/L4 + IDPS)",
    "🌐  WAF protection needed against OWASP Top 10 attacks (L7)",
    "📍  Web App needs the original client IP for geolocation & logging",
    "⚠️  Challenge 1: Azure Firewall SNATs traffic — changes source IP!",
    "⚠️  Challenge 2: VPN users bypass WAF geo-filtering (private IPs can't be geo-resolved)",
    "✅  Solution: App Gateway preserves IP + WAF compound rule secures VPN access",
]
add_bullet_slide_content(slide, items, font_size=20)
add_slide_number(slide, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Diagram
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Architecture Diagram",
                "Internet / VPN → App Gateway (WAF) → Azure Firewall → Web App (Private Endpoint)")

# ─── VNet boundary ───
vnet_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(1.7),
                                     Inches(10.2), Inches(5.4))
vnet_shape.fill.solid()
vnet_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
vnet_shape.line.color.rgb = AZURE_BLUE
vnet_shape.line.width = Pt(2)
add_text_box(slide, 3.0, 1.8, 4.0, 0.3, "Azure Virtual Network  (10.0.0.0/16)",
             font_size=11, bold=True, color=AZURE_BLUE)

# ─── Row 1: Internet flow (top) ───

# Internet Client (outside VNet)
add_shape_box(slide, 0.2, 2.4, 2.0, 1.2, "☁️ Internet\nClient\n(Public IP)", RGBColor(0x60, 0x60, 0x60), WHITE, 11)

# Arrow: Internet → AppGW
add_arrow(slide, 2.3, 2.85, 0.8, 0.35)
add_text_box(slide, 2.3, 2.5, 0.8, 0.3, "HTTP/S", font_size=8, bold=True, color=AZURE_BLUE)

# App Gateway WAF (center-left, spans both rows)
add_shape_box(slide, 3.2, 2.2, 2.4, 2.8,
              "🛡️ App Gateway\nWAF_v2\n\nPublic Frontend\nPrivate Frontend\n\n10.0.0.0/24",
              AZURE_BLUE, WHITE, 10)

# Arrow: AppGW → Firewall
add_arrow(slide, 5.7, 3.3, 0.8, 0.35)
add_text_box(slide, 5.8, 2.9, 0.8, 0.3, "UDR →", font_size=8, bold=True, color=ORANGE)

# Azure Firewall
add_shape_box(slide, 6.6, 2.4, 2.2, 2.2,
              "🔥 Azure Firewall\nPremium\n\nIDPS Inspect\nSNAT\n\n10.0.1.0/26",
              RED, WHITE, 10)

# Arrow: Firewall → Private Endpoint
add_arrow(slide, 8.9, 3.3, 0.8, 0.35)

# Private Endpoint → Web App
add_shape_box(slide, 9.8, 2.4, 2.8, 2.2,
              "🔒 Private Endpoint\n→ Web App\n\nNo Public Access\npublicNetworkAccess:\nDisabled\n\n10.0.3.0/24",
              GREEN, WHITE, 10)

# ─── Row 2: VPN flow (bottom) ───

# VPN Client (outside VNet)
add_shape_box(slide, 0.2, 4.6, 2.0, 1.2, "🔐 VPN Client\n(172.16.0.x)\nP2S VPN", PURPLE, WHITE, 11)

# Arrow: VPN → VPN Gateway
add_arrow(slide, 2.3, 5.05, 0.8, 0.35)
add_text_box(slide, 2.3, 4.7, 0.8, 0.3, "IPsec", font_size=8, bold=True, color=PURPLE)

# VPN Gateway
add_shape_box(slide, 3.2, 5.2, 2.4, 1.2,
              "🌐 VPN Gateway\nP2S · OpenVPN\n10.0.4.0/27",
              PURPLE, WHITE, 10)

# Arrow: VPN GW → AppGW (upward)
shape = slide.shapes.add_shape(MSO_SHAPE.BENT_UP_ARROW, Inches(4.1), Inches(4.7), Inches(0.5), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = PURPLE
shape.line.fill.background()
shape.shadow.inherit = False

# ─── WAF Rules box (below main flow) ───
waf_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(6.0),
                                  Inches(6.5), Inches(0.7))
waf_box.fill.solid()
waf_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
waf_box.line.color.rgb = ORANGE
waf_box.line.width = Pt(1)
tf = waf_box.text_frame
tf.word_wrap = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "WAF: Compound Block Rule (NOT allowed geo AND NOT VPN IP)  |  Managed: OWASP 3.2 + Bot Manager"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = RGBColor(0x80, 0x50, 0x00)
p.font.name = "Segoe UI"

# ─── Log Analytics (right side) ───
add_shape_box(slide, 10.2, 5.0, 2.4, 1.0,
              "📊 Log Analytics\nAccess · WAF · FW",
              ORANGE, WHITE, 10)
add_text_box(slide, 9.0, 5.3, 1.2, 0.3, "← logs", font_size=8, color=ORANGE)

# ─── Flow annotation at bottom ───
add_text_box(slide, 0.3, 6.8, 12.5, 0.4,
             "X-Forwarded-For: Client Public IP → preserved through entire chain → Web App reads original client IP from header ✅",
             font_size=12, bold=True, color=GREEN)
add_slide_number(slide, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Traffic Flow Detail
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Traffic Flow — Step by Step")

steps = [
    "1️⃣  Client sends HTTP/S request to App Gateway Public IP",
    "2️⃣  App Gateway WAF inspects request (Custom Rules → Managed Rules)",
    "3️⃣  App Gateway adds X-Forwarded-For: <Client Public IP> header",
    "4️⃣  App Gateway resolves backend FQDN via Private DNS → PE IP in PE subnet",
    "5️⃣  UDR on AppGW subnet intercepts: route PE subnet → Azure Firewall as next hop",
    "6️⃣  Azure Firewall Premium inspects (IDPS signatures) + SNATs the traffic",
    "7️⃣  Traffic arrives at Web App Private Endpoint",
    "8️⃣  Web App reads original client IP from X-Forwarded-For header ✅",
]
add_bullet_slide_content(slide, steps, font_size=18)
add_slide_number(slide, 4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Key Design Decisions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Key Design Decisions")

items = [
    "🎯  Backend Pool = Web App FQDN (not Firewall IP)",
    "     → Firewall doesn't understand HTTP; health probes would fail",
    "     → UDR silently routes data-plane traffic through Firewall",
    "",
    "🔄  Firewall SNAT forced on private ranges (privateRanges: 255.255.255.255/32)",
    "     → Ensures symmetric routing — return traffic goes back through Firewall",
    "",
    "🔐  Web App: publicNetworkAccess = Disabled",
    "     → Only reachable via Private Endpoint inside the VNet",
    "",
    "🔀  Single UDR on AppGW subnet: PE subnet → Firewall Private IP",
    "     → Works for BOTH internet and VPN traffic — no extra routes needed",
    "",
    "📊  Log Analytics with diagnostic settings on both AppGW and Firewall",
    "     → Full audit trail of every request and its original client IP",
]
add_bullet_slide_content(slide, items, font_size=16)
add_slide_number(slide, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — X-Forwarded-For Explained
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "X-Forwarded-For — How It Works",
                "The HTTP header that preserves the original client IP through proxies")

items = [
    "Client sends request → Source IP: Client Public IP",
    "",
    "App Gateway receives it → adds header:",
    "     X-Forwarded-For: <Client Public IP>",
    "",
    "Azure Firewall SNATs → network source IP changes to Firewall Private IP",
    "     But the HTTP header is untouched! (L4 device, doesn't modify L7 headers)",
    "",
    "Web App receives the request:",
    "     • Remote Address (TCP): Firewall Private IP  ← SNAT'd",
    "     • X-Forwarded-For header: Client Public IP  ← REAL client IP ✅",
    "",
    "⭐ The app reads X-Forwarded-For to get the real client IP for geolocation/logging",
]
add_bullet_slide_content(slide, items, font_size=16)
add_slide_number(slide, 6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Geofencing + VPN: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Use Case 2: Geofencing + VPN — The Problem",
                "When VPN users bypass WAF geo-filtering rules")

items = [
    "🌍  WAF geo-filtering custom rules block traffic from non-allowed countries",
    "     → Works perfectly for internet traffic (public IPs are geo-resolvable)",
    "",
    "🔐  VPN users connect via Point-to-Site VPN → get private IP (e.g. 172.16.0.x)",
    "     → Private IPs CANNOT be geo-resolved → GeoMatch returns 'Unknown' (ZZ)",
    "",
    "⚠️  Problem 1: Geo-block rule sees 'Unknown' ≠ allowed country → BLOCKS VPN users!",
    "⚠️  Problem 2: If VPN traffic bypasses App Gateway → NO WAF protection at all!",
    "",
    "⚠️  Common WRONG approach: Use separate Allow rule for VPN IPs",
    '     → WAF "Allow" action skips ALL subsequent rules INCLUDING managed rules (OWASP)!',
    "     → VPN users would have ZERO protection against SQL injection, XSS, etc.",
    "",
    "❓  How do we allow VPN access while keeping geo-filtering AND managed rules active?",
]
add_bullet_slide_content(slide, items, font_size=16)
add_slide_number(slide, 7)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Geofencing + VPN: The Solution (Compound Rule)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "The Solution — Single Compound Block Rule (Best Practice)",
                "App Gateway private frontend + compound WAF custom rule")

items = [
    "✅  Step 1: Add App Gateway PRIVATE frontend listener",
    "     → VPN clients access App Gateway Private IP — WAF still inspects all traffic!",
    "",
    "✅  Step 2: Create a SINGLE compound Block rule with two AND conditions:",
    "",
    "     Rule: GeoBlockExcludeVpn  (Priority 10, Action: Block)",
    "       Condition 1:  GeoMatch  NOT IN  [allowed country codes]",
    "       Condition 2:  IPMatch   NOT IN  [VPN pool CIDR, internal ranges]",
    "       Both must be TRUE (AND) for the Block to fire",
    "",
    "✅  Step 3: Managed Rules (OWASP 3.2 + Bot Manager) — ALWAYS evaluated",
    "     → When no custom rule fires, traffic falls through to managed rules",
    "     → Protects ALL traffic (internet + VPN) against SQLi, XSS, bots",
    "",
    "⭐ KEY: We do NOT use 'Allow' action — we use a compound 'Block' that excludes VPN",
    "🔗  Ref: learn.microsoft.com/azure/web-application-firewall/ag/geomatch-custom-rules",
]
add_bullet_slide_content(slide, items, font_size=15)
add_slide_number(slide, 8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — WAF Rule Evaluation Flow
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "WAF Rule Evaluation — How Traffic Is Handled",
                "Single compound Block rule + OWASP managed rules")

items = [
    "  Rule: GeoBlockExcludeVpn  (Action: Block)",
    "    Cond 1: GeoMatch NOT IN [allowed countries]   AND",
    "    Cond 2: IPMatch NOT IN [VPN pool, internal IPs]",
    "",
    "  Source                          │ Cond 1          │ Cond 2          │ Rule Fires? │ Managed Rules  │ Result",
    "  ────────────────────── │ ──────── │ ──────── │ ────────│ ──────────── │ ──────",
    "  Internet (allowed country)   │ ❌ FALSE       │ —                 │ NO              │ ✅ Evaluated    │ 200 OK",
    "  Internet (blocked country)   │ ✅ TRUE        │ ✅ TRUE        │ YES → BLOCK │ N/A               │ 403",
    "  VPN client (normal)              │ ✅ TRUE        │ ❌ FALSE       │ NO              │ ✅ Evaluated    │ 200 OK",
    "  VPN client + SQL injection    │ ✅ TRUE        │ ❌ FALSE       │ NO              │ 🚫 OWASP       │ 403",
    "",
    "⭐ KEY INSIGHT: VPN traffic doesn't match Cond 2 → Block rule NEVER fires",
    "   → Traffic falls through to managed rules (OWASP 3.2) → SQLi/XSS still blocked!",
    "",
    '⚠️  WHY NOT separate Allow + Block rules? "Allow" action skips managed rules entirely!',
]
add_bullet_slide_content(slide, items, font_size=14)
add_slide_number(slide, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Architecture Components
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Architecture Components", "Resources deployed via Bicep IaC")

resources = [
    "Component                             │ Purpose",
    "────────────────────────────────── │ ──────────────────────────────────────────",
    "App Gateway (WAF_v2)               │ L7 WAF with public + private frontends, OWASP 3.2 managed rules",
    "WAF Policy + Custom Rule           │ Compound geo-block (exclude VPN IPs), Bot Manager 1.0",
    "Azure Firewall (Premium)           │ L3/L4 inspection, IDPS (signature-based), forced SNAT",
    "Virtual Network (5 subnets)        │ AppGW, Firewall, AppService, PE, Gateway subnets",
    "Web App + Private Endpoint        │ publicNetworkAccess: Disabled — only reachable via PE",
    "Private DNS Zone                       │ privatelink.azurewebsites.net → PE IP resolution",
    "VPN Gateway (P2S)                    │ Point-to-Site, OpenVPN protocol, certificate auth",
    "Route Table + UDR                     │ PE subnet traffic → Firewall as next hop",
    "NSG on PE Subnet                     │ Only allows inbound from Firewall subnet",
    "Log Analytics Workspace            │ Diagnostic logs from AppGW, Firewall, WAF",
]
add_bullet_slide_content(slide, resources, font_size=13, left=0.8)
add_slide_number(slide, 10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Demo: HTTP Echo Proof
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 1: HTTP Echo — X-Forwarded-For Proof",
                "curl http://<App Gateway Public IP> returns all headers the Web App received")

items = [
    "Command:  curl http://<App Gateway Public IP>",
    "",
    "Key headers in the JSON response:",
    "",
    "  x-forwarded-for:    <Client Public IP>:<port>, <AppGW Internal IP>",
    "      └→ First value = YOUR real public IP ✅",
    "      └→ Second value = App Gateway's internal IP",
    "",
    "  x-client-ip:          <AppGW Internal IP>",
    "      └→ After Firewall SNAT — NOT your real IP",
    "",
    "  x-original-host:    <App Gateway Public IP>",
    "  x-appgw-trace-id:  <trace-id> (proves it went through AppGW)",
    "",
    "⭐ Web App can use X-Forwarded-For for geolocation, rate limiting, logging",
]
add_bullet_slide_content(slide, items, font_size=15)
add_slide_number(slide, 11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Demo: WAF Blocks SQL Injection
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 2: WAF Blocks SQL Injection (Managed Rules)",
                "OWASP 3.2 managed rules protect against attacks — internet AND VPN")

items = [
    "Send SQL injection through App Gateway (public):",
    "  curl \"http://<App Gateway Public IP>/?id=1' OR '1'='1\"",
    "",
    "Expected result: ❌ 403 Forbidden — OWASP rule 942100 triggered!",
    "",
    "Send SQL injection from VPN (private frontend):",
    "  curl \"http://<App Gateway Private IP>/?id=1' OR '1'='1\"",
    "",
    "Expected result: ❌ 403 Forbidden — same OWASP rule blocks it!",
    "",
    "✅ Compound Block rule does NOT fire for VPN (Cond 2 = FALSE)",
    "✅ Traffic falls through → managed rules evaluate → SQLi blocked!",
    "",
    "⭐ PROOF: VPN clients bypass geo-filtering but NOT OWASP managed rules!",
]
add_bullet_slide_content(slide, items, font_size=16)
add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Demo: Geo-Filtering Proof (Screenshot)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 3: WAF Compound Rule — Azure Portal Proof",
                "Single GeoBlockExcludeVpn rule with two AND conditions (Block action)")

# Add the screenshot
import os
screenshot_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "waf-custom-rule-screenshot.png")
if os.path.exists(screenshot_path):
    slide.shapes.add_picture(screenshot_path, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.2))
else:
    add_text_box(slide, 1.0, 3.0, 11.0, 1.0,
                 "[Screenshot: Azure Portal → WAF Policy → Custom Rules → GeoBlockExcludeVpn]",
                 font_size=20, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Annotation callouts at the bottom
add_text_box(slide, 0.5, 7.0, 12.5, 0.4,
             "Rule: GeoBlockExcludeVpn  |  Cond 1: GeoMatch IS NOT [2 countries]  |  Cond 2: IPMatch DOES NOT CONTAIN [172.16.0.0/24, 10.0.0.0/8]  |  Then: Deny",
             font_size=11, bold=True, color=AZURE_BLUE, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 13)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Demo: VPN Access + Managed Rules
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 4: VPN Access — Geo-Block Bypassed, Managed Rules Active",
                "Proving VPN clients are allowed but still protected")

items = [
    "Step 1: Connect to P2S VPN (get private IP 172.16.0.x)",
    "  → Azure VPN Client → import profile → connect with client certificate",
    "",
    "Step 2: Access App Gateway private frontend from VPN:",
    "  curl http://<App Gateway Private IP>  → 200 OK ✅",
    "  → Compound Block rule: Cond 2 (NOT VPN IP) = FALSE → rule doesn't fire",
    "  → Managed rules pass → access granted",
    "",
    "Step 3: Send SQL injection from VPN:",
    "  curl \"http://<App Gateway Private IP>/?id=1' OR '1'='1\"  → 403 Forbidden ❌",
    "  → Compound Block rule doesn't fire (VPN IP exempted)",
    "  → BUT managed rules (OWASP 3.2) STILL evaluate → SQLi BLOCKED!",
    "",
    "⭐ PROOF: VPN clients bypass geo-filtering but are FULLY protected by OWASP",
    "⭐ This is the Microsoft-recommended approach for securing VPN + WAF",
]
add_bullet_slide_content(slide, items, font_size=15)
add_slide_number(slide, 14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Demo: Direct Access Blocked
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 5: Web App NOT Accessible Directly",
                "Private Endpoint only — no public access")

items = [
    "Test direct access to the Web App:",
    "  curl https://<webapp-name>.azurewebsites.net  → 403 Forbidden ❌",
    "",
    "Expected result: ❌ Connection refused / 403 Forbidden",
    "  → publicNetworkAccess is set to Disabled",
    "  → No public inbound traffic is allowed",
    "",
    "Test through App Gateway (public):",
    "  curl http://<App Gateway Public IP>  → 200 OK ✅",
    "",
    "Test through App Gateway (private, via VPN):",
    "  curl http://<App Gateway Private IP>  → 200 OK ✅",
    "",
    "⭐ Web App is ONLY reachable through the secure chain (AppGW → Firewall → PE)",
]
add_bullet_slide_content(slide, items, font_size=17)
add_slide_number(slide, 15)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Demo: Log Analytics Proof
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo 6: Log Analytics — Full Audit Trail",
                "Portal → Log Analytics Workspace → Logs")

items = [
    "Query 1 — App Gateway Access Logs (real client IPs):",
    "  AzureDiagnostics | where Category == 'ApplicationGatewayAccessLog'",
    "  | project TimeGenerated, clientIP_s, requestUri_s, httpStatus_d",
    "",
    "Query 2 — WAF Firewall Logs (custom + managed rule actions):",
    "  AzureDiagnostics | where Category == 'ApplicationGatewayFirewallLog'",
    "  | project TimeGenerated, clientIp_s, ruleId_s, action_s, Message",
    "",
    "  Look for: GeoBlockExcludeVpn (custom) + 942100 SQLi (OWASP managed)",
    "",
    "Query 3 — Azure Firewall Network Logs:",
    "  AZFWNetworkRule | project TimeGenerated, SourceIp, DestinationIp, Action",
    "",
    "✅ All three log sources confirm end-to-end traffic inspection",
    "✅ clientIP_s = real client IP for forensic investigation",
]
add_bullet_slide_content(slide, items, font_size=14)
add_slide_number(slide, 16)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Summary
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 1.0, 0.4, 11.0, 0.8, "Summary", font_size=36, bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)

summary_items = [
    "Use Case 1: X-Forwarded-For Preservation",
    "  ✅  App Gateway adds X-Forwarded-For → preserved through Firewall SNAT",
    "  ✅  Web App reads original client IP from header (not network layer)",
    "",
    "Use Case 2: Geofencing + VPN Security",
    "  ✅  Single compound Block rule: Block if (NOT allowed geo) AND (NOT VPN IP)",
    "  ✅  VPN traffic doesn't match → rule doesn't fire → managed rules evaluate",
    "  ✅  Managed rules (OWASP 3.2) protect ALL traffic — internet AND VPN",
    "",
    "Architecture",
    "  ✅  Zero public access to Web App (Private Endpoint only)",
    "  ✅  Azure Firewall Premium with IDPS inspects all traffic",
    "  ✅  UDR + NSG ensure no bypass path exists",
    "  ✅  Full audit trail via Log Analytics",
    "",
    "⚠️  Lesson: Never use separate 'Allow' rule — it skips managed rules!",
    "🔗  learn.microsoft.com/azure/web-application-firewall/ag/geomatch-custom-rules",
]
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.0), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(summary_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.space_after = Pt(5)
    if item.startswith("Use Case") or item.startswith("Architecture"):
        p.font.bold = True
        p.font.size = Pt(19)
    if item.startswith("⚠️"):
        p.font.color.rgb = RGBColor(0xFF, 0xCC, 0x00)
        p.font.bold = True

add_slide_number(slide, 17)


# ── Save ────────────────────────────────────────────────────
output_path = "AppGW-Firewall-WebApp-Demo-westus2.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")

